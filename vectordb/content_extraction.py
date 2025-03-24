import io
import aiohttp
from urllib.parse import urlparse
import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import re

async def parse_file_content(url: str):
    """Parse content from PDF, DOCX, or PPTX file at the given URL."""
    
    # Download file
    async with aiohttp.ClientSession() as session:
        async with session.get(url, headers={"User-Agent": "Mozilla/5.0"}) as response:
            if response.status != 200:
                return f"Error downloading file: {response.status}"
            file_bytes = await response.read()
    
    # Check file signature/magic bytes
    file_type = None
    bytes_data = file_bytes[:8]  # First few bytes for signature detection
    
    # PDF signature: %PDF
    if file_bytes[:4] == b'%PDF':
        file_type = 'pdf'
    # DOCX, PPTX (ZIP-based formats)
    elif bytes_data[:2] == b'PK':
        # Further inspect the ZIP contents for Office XML formats
        byte_stream = io.BytesIO(file_bytes)
        
        # Try to load as PPTX first (since you mentioned this specific URL is a PPTX)
        try:
            Presentation(byte_stream)
            file_type = 'pptx'
        except:
            byte_stream.seek(0)
            try:
                Document(byte_stream)
                file_type = 'docx'
            except:
                # If both fail, check for content markers
                if b'ppt/' in file_bytes[:4000] or b'presentation' in file_bytes[:4000]:
                    file_type = 'pptx'
                elif b'word/' in file_bytes[:4000] or b'document.xml' in file_bytes[:4000]:
                    file_type = 'docx'
    
    # Process based on detected file type
    text = ""
    try:
        if file_type == 'pdf':
            doc = fitz.open(stream=io.BytesIO(file_bytes), filetype="pdf")
            for page in doc:
                text += page.get_text() + "\n\n"
            doc.close()
        elif file_type == 'docx':
            doc = Document(io.BytesIO(file_bytes))
            text = "\n".join([p.text for p in doc.paragraphs if p.text])
        elif file_type == 'pptx':
            prs = Presentation(io.BytesIO(file_bytes))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
                text += "\n"
        else:
            # If still unable to determine, try the most common formats
            try:
                doc = fitz.open(stream=io.BytesIO(file_bytes), filetype="pdf")
                for page in doc:
                    text += page.get_text() + "\n\n"
                doc.close()
                if text.strip():
                    return text
            except:
                pass
                
            try:
                prs = Presentation(io.BytesIO(file_bytes))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text += shape.text + "\n"
                    text += "\n"
                if text.strip():
                    return text
            except:
                pass
                
            try:
                doc = Document(io.BytesIO(file_bytes))
                text = "\n".join([p.text for p in doc.paragraphs if p.text])
                if text.strip():
                    return text
            except:
                text = "Unable to determine or process file type"
    except Exception as e:
        text = f"Error processing file: {str(e)}"
    
    return text

def parse_html_content(html_content: str) -> str:
        """
        Parse HTML content to extract plain text.
        
        Args:
            html_content: HTML content string to parse
            
        Returns:
            Plain text extracted from HTML content
        """
        if not html_content or html_content == "None":
            return ""
        
        try:
            from html.parser import HTMLParser
            
            class HTMLTextExtractor(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.text_parts = []
                    self.in_script = False
                    self.in_style = False
                    
                def handle_starttag(self, tag, attrs):
                    if tag.lower() == "script":
                        self.in_script = True
                    elif tag.lower() == "style":
                        self.in_style = True
                    elif tag.lower() == "br" or tag.lower() == "p":
                        self.text_parts.append("\n")
                    elif tag.lower() == "li":
                        self.text_parts.append("\n• ")
                
                def handle_endtag(self, tag):
                    if tag.lower() == "script":
                        self.in_script = False
                    elif tag.lower() == "style":
                        self.in_style = False
                    elif tag.lower() in ["div", "h1", "h2", "h3", "h4", "h5", "h6", "tr"]:
                        self.text_parts.append("\n")
                
                def handle_data(self, data):
                    if not self.in_script and not self.in_style:
                        # Only append non-empty strings after stripping whitespace
                        text = data.strip()
                        if text:
                            self.text_parts.append(text)
            
                def get_text(self):
                    # Join all text parts and normalize whitespace
                    text = " ".join(self.text_parts)
                    # Replace multiple whitespace with a single space
                    text = re.sub(r'\s+', ' ', text)
                    # Replace multiple newlines with a single newline
                    text = re.sub(r'\n+', '\n', text)
                    return text.strip()
            
            extractor = HTMLTextExtractor()
            extractor.feed(html_content)
            return extractor.get_text()
            
        except Exception as e:
            print(f"Error parsing HTML content: {e}")
            # Fallback to a simple tag stripping approach if the parser fails
            text = re.sub(r'<[^>]*>', ' ', html_content)
            text = re.sub(r'\s+', ' ', text)
            return text.strip()